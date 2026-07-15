#!/usr/bin/env -S uv run --with python-pptx --with lxml python3
"""Render presentation.md into a branded .pptx using a Google Slides template.

Reads the Marp deck (presentation.md), maps each slide onto the template's
native layouts (so the output is fully editable and on-brand), carries the
speaker notes across, and writes a .pptx ready to import into Google Slides.

Usage:
    build_gslides_pptx.py \
        --md presentation/presentation.md \
        --template "SIGGRAPH 2026 template.pptx" \
        --out presentation/presentation-gslides.pptx
"""
from __future__ import annotations

import argparse
import re
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Emu, Pt

ACCENT = RGBColor(0xFF, 0xAB, 0x40)  # theme accent4, echoes the Marp gold

# ---------------------------------------------------------------------------
# Parsing
# ---------------------------------------------------------------------------


@dataclass
class Run:
    text: str
    bold: bool = False
    italic: bool = False
    strike: bool = False


@dataclass
class Para:
    runs: list[Run] = field(default_factory=list)
    level: int = 0
    bullet: bool = False
    fine: bool = False  # small footnote text


@dataclass
class Slide:
    classes: list[str] = field(default_factory=list)
    title_runs: list[list[Run]] = field(default_factory=list)  # heading lines
    body: list[Para] = field(default_factory=list)
    notes: str = ""


INLINE_RE = re.compile(
    r"(\*\*.+?\*\*|__.+?__|<b>.*?</b>|<strong>.*?</strong>"
    r"|\*.+?\*|<em>.*?</em>|<i>.*?</i>)",
    re.S,
)


def parse_inline(text: str) -> list[Run]:
    """Turn a fragment of markdown/HTML into styled runs."""
    text = re.sub(r"!\[[^\]]*\]\([^)]*\)", "", text)  # md images -> drop (no raster)
    text = re.sub(r"<a\b[^>]*>(.*?)</a>", r"\1", text, flags=re.S)  # links -> text
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)  # md links -> text
    text = text.replace("<br>", "\n").replace("<br/>", "\n").replace("<br />", "\n")

    runs: list[Run] = []
    for chunk in INLINE_RE.split(text):
        if not chunk:
            continue
        bold = italic = False
        inner = chunk
        if chunk.startswith(("**", "__")) and chunk.endswith(("**", "__")):
            bold, inner = True, chunk[2:-2]
        elif chunk.startswith(("<b>", "<strong>")):
            bold, inner = True, re.sub(r"</?(b|strong)>", "", chunk)
        elif chunk.startswith(("<em>", "<i>")):
            italic, inner = True, re.sub(r"</?(em|i)>", "", chunk)
        elif chunk.startswith("*") and chunk.endswith("*"):
            italic, inner = True, chunk[1:-1]
        inner = _strip_html(inner)
        if inner:
            runs.append(Run(inner, bold=bold, italic=italic))
    return runs


def _strip_html(s: str) -> str:
    s = re.sub(r"<[^>]+>", "", s)
    return (
        s.replace("&amp;", "&")
        .replace("&lt;", "<")
        .replace("&gt;", ">")
        .replace("&nbsp;", " ")
        .replace("&hellip;", "\u2026")
    )


def parse_slide(block: str) -> Slide:
    sl = Slide()

    # pull out HTML comments; classify directives vs. speaker notes
    for m in re.finditer(r"<!--(.*?)-->", block, re.S):
        inner = m.group(1).strip()
        if inner.startswith("_"):
            cm = re.search(r"_class:\s*(.+)", inner)
            if cm:
                sl.classes = cm.group(1).strip().split()
        else:
            sl.notes = _notes_text(inner)
    block = re.sub(r"<!--.*?-->", "", block, flags=re.S)

    # unwrap block-level HTML paragraphs into their own lines
    ul_items: list[tuple[str, str]] = []  # (li-class, text)

    def _grab_list(m: re.Match) -> str:
        for li in re.finditer(r'<li(?:\s+class="([^"]*)")?>(.*?)</li>', m.group(1), re.S):
            ul_items.append((li.group(1) or "", li.group(2).strip()))
        return "\n\uE000LIST\uE000\n"

    block = re.sub(r"<ul>(.*?)</ul>", _grab_list, block, flags=re.S)

    lines = [ln.rstrip() for ln in block.strip().splitlines()]
    li_iter = iter(ul_items)

    for raw in lines:
        line = raw.strip()
        if not line:
            continue

        # standalone markdown image -> skip (python-pptx can't embed SVG)
        if re.fullmatch(r"!\[[^\]]*\]\([^)]*\)", line):
            continue

        if line == "\uE000LIST\uE000":
            for cls, txt in li_iter:
                runs = parse_inline(txt)
                if "fail" in cls:
                    for r in runs:
                        r.strike = True
                sl.body.append(Para(runs=runs, level=0, bullet=True))
            continue

        # headings -> title lines
        hm = re.match(r"(#{1,6})\s+(.*)", line)
        if hm:
            sl.title_runs.append(parse_inline(hm.group(2)))
            continue

        # blockquote -> title (used by quote slides)
        if line.startswith(">"):
            sl.title_runs.append(parse_inline(line.lstrip("> ").strip()))
            continue

        # markdown bullet
        bm = re.match(r"[-*]\s+(.*)", line)
        if bm:
            sl.body.append(Para(runs=parse_inline(bm.group(1)), bullet=True))
            continue

        # <p class="fine|lede|who|...">
        pm = re.match(r'<p(?:\s+class="([^"]*)")?>(.*?)</p>', line, re.S)
        if pm:
            cls = pm.group(1) or ""
            if "placeholder" in cls:
                continue  # drop the 🚧 art placeholders
            sl.body.append(Para(runs=parse_inline(pm.group(2)), fine="fine" in cls))
            continue

        # bare paragraph text
        sl.body.append(Para(runs=parse_inline(line)))

    return sl


def _notes_text(inner: str) -> str:
    out = []
    for ln in inner.splitlines():
        ln = ln.strip()
        if ln.startswith("- "):
            ln = ln[2:]
        out.append(_strip_html(ln.replace("‖", "  //  ")))
    return "\n".join(x for x in out if x)


def parse_deck(md_path: Path) -> list[Slide]:
    text = md_path.read_text()
    # drop YAML front matter and the <style> block
    text = re.sub(r"^---\n.*?\n---\n", "", text, count=1, flags=re.S)
    text = re.sub(r"<style>.*?</style>", "", text, flags=re.S)
    blocks = re.split(r"\n---\n", text)
    return [parse_slide(b) for b in blocks if b.strip()]


# ---------------------------------------------------------------------------
# Rendering
# ---------------------------------------------------------------------------


def _find_layout(prs: Presentation, name: str):
    for lay in prs.slide_layouts:
        if lay.name == name:
            return lay
    raise KeyError(name)


def _ph(slide, *types):
    """First placeholder matching any of the given type names."""
    for ph in slide.placeholders:
        if str(ph.placeholder_format.type).split(" ")[0] in types:
            yield ph


def _set_runs(tf, paras: list[Para], *, title=False):
    tf.clear()
    first = True
    for para in paras:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = para.level
        # split runs on embedded newlines into separate paragraphs
        segments = _split_newlines(para.runs)
        for si, seg in enumerate(segments):
            if si > 0:
                p = tf.add_paragraph()
                p.level = para.level
            for r in seg:
                run = p.add_run()
                run.text = r.text
                if r.bold:
                    run.font.bold = True
                    run.font.color.rgb = ACCENT
                if r.italic:
                    run.font.italic = True
                if para.fine:
                    run.font.size = Pt(12)


def _split_newlines(runs: list[Run]) -> list[list[Run]]:
    segments: list[list[Run]] = [[]]
    for r in runs:
        parts = r.text.split("\n")
        for i, part in enumerate(parts):
            if i > 0:
                segments.append([])
            if part:
                segments[-1].append(Run(part, r.bold, r.italic, r.strike))
    return segments


def _title_text(tf, title_runs: list[list[Run]]):
    paras = [Para(runs=line) for line in title_runs]
    _set_runs(tf, paras, title=True)


def _fit(tf):
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass


def _place(ph, *, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    ph.left, ph.top, ph.width, ph.height = (
        Emu(int(left * 914400)),
        Emu(int(top * 914400)),
        Emu(int(width * 914400)),
        Emu(int(height * 914400)),
    )
    ph.text_frame.vertical_anchor = anchor


def _layout_dark(title_ph, sub_ph):
    """Reposition a dark title/section slide so title and subtitle never collide."""
    _place(title_ph, left=1.37, top=2.1, width=16.0, height=3.9)
    _fit(title_ph.text_frame)
    if sub_ph is not None:
        _place(sub_ph, left=1.37, top=6.4, width=14.5, height=3.0)
        _fit(sub_ph.text_frame)


def _delete_slide(prs: Presentation, index: int):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[index].get(
        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
    )
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[index])


def _set_notes(slide, text: str):
    if not text:
        return
    slide.notes_slide.notes_text_frame.text = text


def build(md_path: Path, template_path: Path, out_path: Path, *, keep_cover=True):
    slides = parse_deck(md_path)
    prs = Presentation(str(template_path))

    # capture the content layout from the template's own sample content slide
    content_layout = prs.slides[1].slide_layout
    title_layout = _find_layout(prs, "TITLE")
    section_layout = _find_layout(prs, "SECTION_HEADER")

    # strip the template's sample slides, optionally keeping the branded cover
    start = len(prs.slides._sldIdLst)
    for i in reversed(range(start)):
        if keep_cover and i == 0:
            continue
        _delete_slide(prs, i)

    for sl in slides:
        cls = set(sl.classes)
        if "title" in cls:
            slide = prs.slides.add_slide(title_layout)
            title_ph = next(_ph(slide, "TITLE"), None)
            if title_ph is not None:
                _title_text(title_ph.text_frame, sl.title_runs)
            sub_ph = next(_ph(slide, "SUBTITLE", "BODY"), None)
            if sub_ph is not None and sl.body:
                _set_runs(sub_ph.text_frame, sl.body)
            _layout_dark(title_ph, sub_ph if sl.body else None)
        elif cls & {"section", "spine", "quote"}:
            slide = prs.slides.add_slide(section_layout)
            title_ph = next(_ph(slide, "TITLE"), None)
            if title_ph is not None:
                _title_text(title_ph.text_frame, sl.title_runs)
            sub_ph = next(_ph(slide, "SUBTITLE", "BODY"), None)
            if sub_ph is not None and sl.body:
                _set_runs(sub_ph.text_frame, sl.body)
            _layout_dark(title_ph, sub_ph if sl.body else None)
        else:
            slide = prs.slides.add_slide(content_layout)
            title_ph = next(_ph(slide, "TITLE"), None)
            if title_ph is not None:
                _title_text(title_ph.text_frame, sl.title_runs)
                _place(title_ph, left=0.84, top=0.55, width=15.5, height=1.5)
                _fit(title_ph.text_frame)
            bodies = list(_ph(slide, "BODY", "SUBTITLE"))
            if bodies:
                _set_runs(bodies[0].text_frame, sl.body)
        _set_notes(slide, sl.notes)

    prs.save(str(out_path))
    return len(slides)


def main():
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--md", required=True, type=Path)
    ap.add_argument("--template", required=True, type=Path)
    ap.add_argument("--out", required=True, type=Path)
    ap.add_argument("--drop-cover", action="store_true",
                    help="omit the template's branded cover slide")
    args = ap.parse_args()
    n = build(args.md, args.template, args.out, keep_cover=not args.drop_cover)
    print(f"wrote {args.out} ({n} content slides + cover)")


if __name__ == "__main__":
    main()
